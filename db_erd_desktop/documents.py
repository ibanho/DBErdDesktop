from __future__ import annotations

import html
import textwrap
import zipfile
from datetime import datetime
from pathlib import Path

from .models import DatabaseModel, TableModel


def export_document(
    model: DatabaseModel,
    output_path: Path,
    logical_png: Path,
    physical_png: Path,
) -> None:
    suffix = output_path.suffix.lower()
    if suffix == ".docx":
        _export_docx(model, output_path, logical_png, physical_png)
    elif suffix == ".pptx":
        _export_pptx(model, output_path, logical_png, physical_png)
    elif suffix == ".hwpx":
        _export_hwpx(model, output_path, logical_png, physical_png)
    else:
        raise ValueError("지원 문서 형식은 .docx, .pptx, .hwpx 입니다.")


def _export_docx(model: DatabaseModel, output_path: Path, logical_png: Path, physical_png: Path) -> None:
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError as exc:
        raise RuntimeError("DOCX 출력을 위해 python-docx가 필요합니다.") from exc

    document = Document()
    document.add_heading("Database ERD Documentation", 0)
    document.add_paragraph(f"DBMS: {model.dbms}")
    document.add_paragraph(f"Database/Schema: {model.namespace}")
    document.add_paragraph(f"Generated: {datetime.now():%Y-%m-%d %H:%M}")

    document.add_heading("Logical ERD", level=1)
    document.add_picture(str(logical_png), width=Inches(6.6))
    document.add_heading("Physical ERD", level=1)
    document.add_picture(str(physical_png), width=Inches(6.6))

    document.add_heading("Tables", level=1)
    for table in model.tables:
        document.add_heading(f"{table.logical_name} ({table.qualified_name})", level=2)
        if table.comment:
            document.add_paragraph(table.comment)
        grid = document.add_table(rows=1, cols=7)
        grid.style = "Table Grid"
        headers = ["Column", "Logical Name", "Type", "PK", "FK", "Nullable", "Default"]
        for index, header in enumerate(headers):
            grid.rows[0].cells[index].text = header
        for column in table.columns:
            cells = grid.add_row().cells
            values = [
                column.name,
                column.logical_name,
                column.data_type,
                "Y" if column.is_primary_key else "",
                column.foreign_key or "",
                "Y" if column.nullable else "N",
                column.default or "",
            ]
            for index, value in enumerate(values):
                cells[index].text = value

    document.add_heading("Relationships", level=1)
    if model.relationships:
        for relationship in model.relationships:
            document.add_paragraph(
                f"{relationship.from_table}({', '.join(relationship.from_columns)}) -> "
                f"{relationship.to_table}({', '.join(relationship.to_columns)})",
                style="List Bullet",
            )
    else:
        document.add_paragraph("No foreign key relationships found.")

    document.save(str(output_path))


def _export_pptx(model: DatabaseModel, output_path: Path, logical_png: Path, physical_png: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("PPTX 출력을 위해 python-pptx가 필요합니다.") from exc

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Database ERD Documentation"
    title_slide.placeholders[1].text = (
        f"{model.dbms} / {model.namespace}\nGenerated: {datetime.now():%Y-%m-%d %H:%M}"
    )

    for title, image in [("Logical ERD", logical_png), ("Physical ERD", physical_png)]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(str(image), Inches(0.45), Inches(1.0), width=Inches(9.1))

    for table in model.tables:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{table.logical_name} ({table.name})"
        rows = min(len(table.columns) + 1, 15)
        table_shape = slide.shapes.add_table(rows, 5, Inches(0.4), Inches(1.1), Inches(9.2), Inches(5.6))
        ppt_table = table_shape.table
        headers = ["Column", "Logical", "Type", "Key", "Nullable"]
        for col_index, header in enumerate(headers):
            ppt_table.cell(0, col_index).text = header
        for row_index, column in enumerate(table.columns[: rows - 1], start=1):
            key = "/".join(part for part in ["PK" if column.is_primary_key else "", "FK" if column.foreign_key else ""] if part)
            values = [column.name, column.logical_name, column.data_type, key, "Y" if column.nullable else "N"]
            for col_index, value in enumerate(values):
                cell = ppt_table.cell(row_index, col_index)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)

    presentation.save(str(output_path))


def _export_hwpx(model: DatabaseModel, output_path: Path, logical_png: Path, physical_png: Path) -> None:
    paragraphs = _document_lines(model, logical_png, physical_png)
    section = _hwpx_section(paragraphs)
    preview = "\n".join(paragraphs)

    with zipfile.ZipFile(output_path, "w") as archive:
        archive.writestr("mimetype", "application/hwp+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr("version.xml", _VERSION_XML)
        archive.writestr("META-INF/manifest.xml", _manifest_xml())
        archive.writestr("Contents/content.hpf", _content_hpf())
        archive.writestr("Contents/header.xml", _HEADER_XML)
        archive.writestr("Contents/section0.xml", section)
        archive.writestr("Preview/PrvText.txt", preview)


def _document_lines(model: DatabaseModel, logical_png: Path, physical_png: Path) -> list[str]:
    lines = [
        "Database ERD Documentation",
        f"DBMS: {model.dbms}",
        f"Database/Schema: {model.namespace}",
        f"Generated: {datetime.now():%Y-%m-%d %H:%M}",
        "",
        "ERD Images",
        f"Logical ERD PNG: {logical_png.name}",
        f"Physical ERD PNG: {physical_png.name}",
        "",
        "Tables",
    ]
    for table in model.tables:
        lines.extend(_table_lines(table))
    lines.append("")
    lines.append("Relationships")
    if model.relationships:
        for relationship in model.relationships:
            lines.append(
                f"- {relationship.from_table}({', '.join(relationship.from_columns)}) -> "
                f"{relationship.to_table}({', '.join(relationship.to_columns)})"
            )
    else:
        lines.append("- No foreign key relationships found.")
    return lines


def _table_lines(table: TableModel) -> list[str]:
    lines = ["", f"[{table.logical_name}] {table.qualified_name}"]
    if table.comment:
        lines.append(table.comment)
    for column in table.columns:
        key = " ".join(part for part in ["PK" if column.is_primary_key else "", "FK" if column.foreign_key else ""] if part)
        key = f" {key}" if key else ""
        nullable = "NULL" if column.nullable else "NOT NULL"
        lines.append(f"- {column.name} ({column.logical_name}) {column.data_type} {nullable}{key}")
    return lines


def _hwpx_section(lines: list[str]) -> str:
    body = []
    next_id = 1000000001
    body.append(_section_paragraph(next_id))
    next_id += 1
    for line in lines:
        body.append(_paragraph(next_id, line))
        next_id += 1
    return (
        "<?xml version='1.0' encoding='UTF-8'?>\n"
        '<hs:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph" '
        'xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section">\n'
        + "\n".join(body)
        + "\n</hs:sec>\n"
    )


def _section_paragraph(paragraph_id: int) -> str:
    return textwrap.dedent(
        f"""\
        <hp:p id="{paragraph_id}" paraPrIDRef="0" styleIDRef="0" pageBreak="0" columnBreak="0" merged="0">
          <hp:run charPrIDRef="0">
            <hp:secPr id="" textDirection="HORIZONTAL" spaceColumns="0" tabStop="8000" tabStopVal="LEFT" tabStopUnit="HWPUNIT" outlineShapeIDRef="0" memoShapeIDRef="0">
              <hp:grid lineGrid="0" charGrid="0" wonggojiFormat="0"/>
              <hp:startNum pageStartsOn="BOTH" page="0" pic="0" tbl="0" equation="0"/>
              <hp:visibility hideFirstHeader="0" hideFirstFooter="0" hideFirstMasterPage="0" border="SHOW_ALL" fill="SHOW_ALL" hideFirstPageNum="0" hideFirstEmptyLine="0" showLineNumber="0"/>
              <hp:lineNumberShape textShape="DIGIT" position="LEFT" countBy="0" distance="0" startNumber="1"/>
              <hp:pagePr landscape="0" width="59528" height="84186" gutterType="LEFT_ONLY">
                <hp:margin header="4252" footer="4252" gutter="0" left="8504" right="8504" top="5668" bottom="4252"/>
              </hp:pagePr>
              <hp:footNotePr/>
              <hp:endNotePr/>
              <hp:pageBorderFill type="BOTH" borderFillIDRef="1" textBorder="PAPER" headerInside="0" footerInside="0" fillArea="PAPER"/>
              <hp:masterPage idRef="0"/>
            </hp:secPr>
            <hp:ctrl><hp:colPr id="" type="NEWSPAPER" layout="LEFT" colCount="1" sameSz="1" sameGap="0"/></hp:ctrl>
          </hp:run>
          <hp:run charPrIDRef="0"><hp:t/></hp:run>
        </hp:p>"""
    )


def _paragraph(paragraph_id: int, text: str) -> str:
    escaped = html.escape(text)
    text_node = f"<hp:t>{escaped}</hp:t>" if escaped else "<hp:t/>"
    return (
        f'<hp:p id="{paragraph_id}" paraPrIDRef="0" styleIDRef="0" pageBreak="0" columnBreak="0" merged="0">'
        f'<hp:run charPrIDRef="0">{text_node}</hp:run></hp:p>'
    )


def _manifest_xml() -> str:
    entries = [
        ("version.xml", "text/xml"),
        ("Contents/header.xml", "text/xml"),
        ("Contents/section0.xml", "text/xml"),
        ("Contents/content.hpf", "text/xml"),
        ("Preview/PrvText.txt", "text/plain"),
    ]
    items = "\n".join(
        f'  <manifest:file-entry manifest:media-type="{media_type}" manifest:full-path="{path}"/>'
        for path, media_type in entries
    )
    return (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">\n'
        '  <manifest:file-entry manifest:media-type="application/hwp+zip" manifest:full-path="/"/>\n'
        f"{items}\n"
        "</manifest:manifest>\n"
    )


def _content_hpf() -> str:
    return textwrap.dedent(
        """\
        <?xml version="1.0" encoding="UTF-8"?>
        <opf:package xmlns:opf="http://www.idpf.org/2007/opf" version="1.0" unique-identifier="uid">
          <opf:metadata>
            <opf:title>Database ERD Documentation</opf:title>
            <opf:language>ko-KR</opf:language>
          </opf:metadata>
          <opf:manifest>
            <opf:item id="header" href="header.xml" media-type="text/xml"/>
            <opf:item id="section0" href="section0.xml" media-type="text/xml"/>
          </opf:manifest>
          <opf:spine>
            <opf:itemref idref="section0"/>
          </opf:spine>
        </opf:package>
        """
    )


_VERSION_XML = """<?xml version="1.0" encoding="UTF-8"?><HCFVersion version="1.0"/>"""

_HEADER_XML = textwrap.dedent(
    """\
    <?xml version="1.0" encoding="UTF-8"?>
    <hh:head xmlns:hh="http://www.hancom.co.kr/hwpml/2011/head" xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph" version="1.0">
      <hh:beginNum page="1" footnote="1" endnote="1" pic="1" tbl="1" equation="1"/>
      <hh:refList>
        <hh:fontfaces itemCnt="1">
          <hh:fontface lang="HANGUL" itemCnt="1"><hh:font id="0" face="함초롬바탕" type="TTF"/></hh:fontface>
        </hh:fontfaces>
        <hh:borderFills itemCnt="2">
          <hh:borderFill id="1" threeD="0" shadow="0" centerLine="NONE"><hh:slash type="NONE" Crooked="0" isCounter="0"/><hh:backSlash type="NONE" Crooked="0" isCounter="0"/><hh:leftBorder type="NONE" width="0.1 mm" color="#000000"/><hh:rightBorder type="NONE" width="0.1 mm" color="#000000"/><hh:topBorder type="NONE" width="0.1 mm" color="#000000"/><hh:bottomBorder type="NONE" width="0.1 mm" color="#000000"/><hh:diagonal type="NONE" width="0.1 mm" color="#000000"/></hh:borderFill>
          <hh:borderFill id="2" threeD="0" shadow="0" centerLine="NONE"><hh:slash type="NONE" Crooked="0" isCounter="0"/><hh:backSlash type="NONE" Crooked="0" isCounter="0"/><hh:leftBorder type="NONE" width="0.1 mm" color="#000000"/><hh:rightBorder type="NONE" width="0.1 mm" color="#000000"/><hh:topBorder type="NONE" width="0.1 mm" color="#000000"/><hh:bottomBorder type="NONE" width="0.1 mm" color="#000000"/><hh:diagonal type="NONE" width="0.1 mm" color="#000000"/></hh:borderFill>
        </hh:borderFills>
        <hh:charProperties itemCnt="1">
          <hh:charPr id="0" height="1000" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
            <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
            <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
            <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
            <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
            <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
          </hh:charPr>
        </hh:charProperties>
        <hh:paraProperties itemCnt="1">
          <hh:paraPr id="0" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
            <hh:align horizontal="LEFT" vertical="BASELINE"/>
            <hh:heading type="NONE" idRef="0" level="0"/>
            <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="0" keepLines="0" pageBreakBefore="0" lineWrap="BREAK"/>
            <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
            <hh:margin><hh:intent value="0"/><hh:left value="0"/><hh:right value="0"/><hh:prev value="0"/><hh:next value="0"/></hh:margin>
            <hh:lineSpacing type="PERCENT" value="160"/>
            <hh:border borderFillIDRef="2" offsetLeft="0" offsetRight="0" offsetTop="0" offsetBottom="0" connect="0" ignoreMargin="0"/>
          </hh:paraPr>
        </hh:paraProperties>
        <hh:styles itemCnt="1"><hh:style id="0" type="PARA" name="바탕글" engName="Normal" paraPrIDRef="0" charPrIDRef="0" nextStyleIDRef="0" langID="1042" lockForm="0"/></hh:styles>
        <hh:tabProperties itemCnt="1"><hh:tabPr id="0" autoTabLeft="0" autoTabRight="0"/></hh:tabProperties>
      </hh:refList>
    </hh:head>
    """
)
